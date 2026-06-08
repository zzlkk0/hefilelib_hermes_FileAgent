"""
File Manager MCP Server — ingest module
Extracts text content from various file types.
"""

import os
import json
import subprocess
from pathlib import Path
from typing import Optional

# Try importing optional deps
try:
    import pymupdf
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import jieba
    HAS_JIEBA = True
except ImportError:
    HAS_JIEBA = False


TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".log", ".py", ".js", ".ts", ".sh", ".bash", ".css", ".toml", ".ini", ".cfg"}
IPYNB_EXTENSIONS = {".ipynb"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tiff"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx", ".doc"}
PPTX_EXTENSIONS = {".pptx", ".ppt"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".flac", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}


def extract_text(filepath: str) -> dict:
    """Extract text content from a file. Returns {text, keywords, file_type, error}."""
    path = Path(filepath)
    ext = path.suffix.lower()

    result = {
        "filepath": str(path),
        "filename": path.name,
        "file_type": ext,
        "text": "",
        "keywords": [],
        "error": None
    }

    try:
        if not path.exists():
            result["error"] = f"File not found: {filepath}"
            return result

        # Text files — direct read
        if ext in TEXT_EXTENSIONS:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                result["text"] = f.read(50000)
            result["file_type"] = "text"

        # Jupyter notebooks — parse JSON, extract code + markdown
        elif ext in IPYNB_EXTENSIONS:
            result = _extract_ipynb(path, result)

        # PDF files
        elif ext in PDF_EXTENSIONS:
            result = _extract_pdf(path, result)

        # Word documents
        elif ext in DOCX_EXTENSIONS:
            result = _extract_docx(path, result)

        # PowerPoint
        elif ext in PPTX_EXTENSIONS:
            result = _extract_pptx(path, result)

        # Images — no text extraction, metadata only
        elif ext in IMAGE_EXTENSIONS:
            result["file_type"] = "image"
            result["text"] = f"[Image: {path.name}]"

        # Audio — no STT by default (needs whisper setup)
        elif ext in AUDIO_EXTENSIONS:
            result["file_type"] = "audio"
            result["text"] = f"[Audio: {path.name}]"

        # Video — no processing by default
        elif ext in VIDEO_EXTENSIONS:
            result["file_type"] = "video"
            result["text"] = f"[Video: {path.name}]"

        else:
            result["file_type"] = "unknown"
            result["text"] = f"[Unknown type: {path.name}]"

        # Extract keywords
        if result["text"] and not result["text"].startswith("["):
            result["keywords"] = _extract_keywords(result["text"])

    except Exception as e:
        result["error"] = str(e)

    return result


def _extract_pdf(path: Path, result: dict) -> dict:
    """Extract text from PDF using pymupdf."""
    if not HAS_PYMUPDF:
        result["error"] = "pymupdf not installed. Run: pip install pymupdf pymupdf4llm"
        return result

    try:
        doc = pymupdf.open(str(path))
        text_parts = []
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text)
        doc.close()

        full_text = "\n".join(text_parts)

        # If almost no text extracted, likely scanned PDF
        if len(full_text.strip()) < 50:
            result["text"] = full_text
            result["file_type"] = "scanned_pdf"
            result["error"] = "Low text content — may be scanned PDF. Install marker-pdf for OCR."
        else:
            result["text"] = full_text[:50000]
            result["file_type"] = "pdf"
    except Exception as e:
        result["error"] = f"PDF extraction failed: {e}"

    return result


def _extract_docx(path: Path, result: dict) -> dict:
    """Extract text from Word documents."""
    if not HAS_DOCX:
        result["error"] = "python-docx not installed. Run: pip install python-docx"
        return result

    try:
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result["text"] = "\n".join(paragraphs)[:50000]
        result["file_type"] = "docx"
    except Exception as e:
        result["error"] = f"DOCX extraction failed: {e}"

    return result


def _extract_pptx(path: Path, result: dict) -> dict:
    """Extract text from PowerPoint."""
    if not HAS_PPTX:
        result["error"] = "python-pptx not installed. Run: pip install python-pptx"
        return result

    try:
        prs = Presentation(str(path))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_text.append(t)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        result["text"] = "\n\n".join(text_parts)[:50000]
        result["file_type"] = "pptx"
    except Exception as e:
        result["error"] = f"PPTX extraction failed: {e}"

    return result


def _extract_keywords(text: str, top_k: int = 20) -> list[str]:
    """Extract keywords using jieba (Chinese) or simple word frequency (English)."""
    if not HAS_JIEBA:
        # Fallback: simple space-split + frequency
        words = text.lower().split()
        freq = {}
        for w in words:
            w = w.strip(".,!?;:()[]{}<>\"'，。！？；：（）【】《》").strip()
            if len(w) >= 2:
                freq[w] = freq.get(w, 0) + 1
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [w for w, _ in sorted_words[:top_k]]

    # Use jieba for Chinese-aware segmentation
    import jieba.analyse
    keywords = jieba.analyse.extract_tags(text, topK=top_k, withWeight=False)
    return keywords


def extract_batch(directory: str) -> list[dict]:
    """Extract text from all files in a directory."""
    results = []
    dir_path = Path(directory)
    if not dir_path.exists():
        return [{"error": f"Directory not found: {directory}"}]

    for filepath in dir_path.iterdir():
        if filepath.is_file() and not filepath.name.startswith("."):
            results.append(extract_text(str(filepath)))

    return results


def _extract_ipynb(path: Path, result: dict) -> dict:
    """Extract code and markdown content from a Jupyter notebook."""
    import json
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)
    except (json.JSONDecodeError, Exception) as e:
        result["error"] = f"Failed to parse notebook: {e}"
        return result

    cells = nb.get("cells", [])
    parts = []
    code_count = 0
    md_count = 0

    for i, cell in enumerate(cells):
        cell_type = cell.get("cell_type", "")
        source = cell.get("source", [])

        # source can be a list of strings or a single string
        if isinstance(source, list):
            source_text = "".join(source)
        else:
            source_text = str(source)

        if not source_text.strip():
            continue

        if cell_type == "code":
            code_count += 1
            parts.append(f"# --- Code Cell {code_count} ---\n{source_text}")
        elif cell_type == "markdown":
            md_count += 1
            parts.append(f"# --- Markdown Cell {md_count} ---\n{source_text}")

    result["text"] = "\n\n".join(parts)[:50000]
    result["file_type"] = "ipynb"

    if not result["text"]:
        result["text"] = f"[Jupyter Notebook: {code_count} code cells, {md_count} markdown cells]"

    return result
